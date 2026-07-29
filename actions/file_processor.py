import os
import json
import traceback
from pathlib import Path


def _search_file(name_hint: str) -> str | None:
    name_lower = name_hint.lower().strip() if name_hint else ""
    base = Path(__file__).resolve().parent.parent
    dirs = [
        Path.home() / "Desktop",
        Path.home() / "Documents",
        Path.home() / "Downloads",
        base,
    ]
    for d in [base / "ONYX-IA-main", base / "subidas", base / "uploads"]:
        if d.exists():
            dirs.append(d)
    exclude_dirs = {".venv", "__pycache__", ".git", "node_modules", "venv", "env"}
    priority_exts = {".docx", ".doc", ".pdf", ".xlsx", ".pptx", ".txt"}
    candidates = []
    for d in dirs:
        if d.exists():
            for f in d.rglob(f"*"):
                if f.is_file():
                    parts = f.relative_to(d).parts if f.is_relative_to(d) else f.parts
                    if any(ex in parts for ex in exclude_dirs):
                        continue
                    fn = f.stem.lower()
                    ext = f.suffix.lower()
                    if ext in priority_exts:
                        if not name_lower or name_lower in fn:
                            candidates.append(f)
    if candidates:
        def sort_key(f):
            ext = f.suffix.lower()
            prio = 0 if ext in (".docx", ".doc", ".pdf") else 1
            return (prio, -f.stat().st_mtime)
        candidates.sort(key=sort_key)
        return str(candidates[0])
    return None


def file_processor(parameters: dict, player=None, speak=None) -> str:
    action = parameters.get("action", "").lower()
    file_path = parameters.get("file_path", "")
    instruction = parameters.get("instruction", "")
    fmt = parameters.get("format", "")

    if player and player.current_file and not file_path:
        file_path = player.current_file

    if not file_path:
        found = _search_file("")
        if found:
            file_path = found
        else:
            return "Error: No se especificó la ruta del archivo."

    path = Path(file_path)
    if not path.exists():
        found = _search_file(path.name)
        if found:
            path = Path(found)
            file_path = found
        else:
            if player:
                player.write_log(f"❌ Archivo no encontrado: {file_path}")
            return f"Error: No se encontró el archivo '{path.name}'."

    try:
        ext = path.suffix.lower()

        if ext in (".docx", ".doc"):
            return _process_word(path, action, instruction, player)

        elif ext == ".pdf":
            return _process_pdf(path, action, instruction, player)

        elif ext in (".txt", ".md", ".json", ".xml", ".csv", ".log", ".py", ".js", ".html", ".css"):
            return _process_text(path, action, instruction, fmt, player)

        elif ext in (".xlsx", ".xls"):
            return _process_excel(path, action, instruction, player)

        elif ext in (".pptx", ".ppt"):
            return _process_powerpoint(path, action, player)

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
            return _process_image(path, action, player)

        elif ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac"):
            return "Formato de audio detectado. Usa la herramienta screen_vision o el procesador de audio."

        else:
            text = path.read_text(encoding="utf-8", errors="replace")[:2000]
            return f"Contenido de '{path.name}':\n\n{text}"

    except Exception as e:
        return f"Error al procesar '{path.name}': {str(e)}"


def _process_word(path: Path, action: str, instruction: str, player) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n".join(paragraphs)

        if action in ("resumir", "summarize", "leer", "read") or not action:
            return f"Contenido de '{path.name}':\n\n{full_text[:2000]}"

        elif action in ("extraer_texto", "extract_text"):
            return full_text[:3000]

        elif action in ("info", "informacion"):
            info_lines = [
                f"Archivo: {path.name}",
                f"Tamaño: {path.stat().st_size:,} bytes",
                f"Párrafos: {len(paragraphs)}",
                f"Primeras líneas:\n{full_text[:500]}"
            ]
            return "\n".join(info_lines)

        elif action == "convertir" or action == "convert":
            if player:
                player.write_log(f"📄 Convirtiendo '{path.name}'...")
            return full_text[:3000]

        return f"Contenido de '{path.name}':\n\n{full_text[:2000]}"

    except ImportError:
        return "Error: python-docx no está instalado. Ejecutá: pip install python-docx"
    except Exception as e:
        return f"Error al leer Word: {str(e)}"


def _process_pdf(path: Path, action: str, instruction: str, player) -> str:
    try:
        import PyPDF2
        text = ""
        with open(str(path), "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        if not text.strip():
            return f"El PDF '{path.name}' parece estar escaneado (sin texto seleccionable). Usá la acción 'open' para verlo."
        return f"Contenido de '{path.name}':\n\n{text[:2000]}"
    except ImportError:
        return "Error: PyPDF2 no está instalado."
    except Exception as e:
        return f"Error al leer PDF: {str(e)}"


def _process_text(path: Path, action: str, instruction: str, fmt: str, player) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")

    if action in ("resumir", "summarize"):
        lines = text.split("\n")
        preview = "\n".join(lines[:20])
        return f"Resumen de '{path.name}' ({len(lines)} líneas, {len(text)} caracteres):\n\n{preview}"

    elif action in ("leer", "read") or not action:
        return f"Contenido de '{path.name}':\n\n{text[:2000]}"

    elif action in ("info", "informacion"):
        lines = text.split("\n")
        return f"Archivo: {path.name}\nLíneas: {len(lines)}\nCaracteres: {len(text)}\nPrimeras líneas:\n{text[:500]}"

    return f"Contenido de '{path.name}':\n\n{text[:2000]}"


def _process_excel(path: Path, action: str, instruction: str, player) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        ws = wb.active
        rows = []
        for i, row in enumerate(ws.iter_row()):
            if i > 30:
                break
            rows.append(" | ".join(str(cell.value or "") for cell in row))
        wb.close()
        return f"Contenido de '{path.name}':\n\n" + "\n".join(rows)
    except ImportError:
        return "Error: openpyxl no está instalado."
    except Exception as e:
        return f"Error al leer Excel: {str(e)}"


def _process_powerpoint(path: Path, action: str, player) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        preview = "\n".join(texts[:50])
        return f"Contenido de '{path.name}':\n\n{preview[:5000]}"
    except ImportError:
        return "Error: python-pptx no está instalado."
    except Exception as e:
        return f"Error al leer PowerPoint: {str(e)}"


def _process_image(path: Path, action: str, player) -> str:
    if action in ("abrir", "open", "mostrar", "show") or not action:
        os.startfile(str(path))
        return f"He abierto la imagen '{path.name}'."
    return f"Imagen '{path.name}' detectada. Usá screen_vision para analizarla con IA."
